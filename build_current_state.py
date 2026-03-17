#!/usr/bin/env python3
"""
Build GMA Current State HTML from PowerPoint presentation.
Extracts text, tables, AND images with layout-aware positioning.
Renders with the same design system as the gap analysis dashboard.
"""

import sys, io, os, re, base64
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400
SIDEBAR_THRESHOLD = 7.8 * EMU_PER_INCH  # shapes with left > 7.8" are sidebar
TITLE_THRESHOLD = 1.8 * EMU_PER_INCH    # shapes with top < 1.8" are title area
SOURCE_THRESHOLD = 6.3 * EMU_PER_INCH   # shapes with top > 6.3" may be source

NOISE_TEXTS = {
    'AUTHORIZED SIGNERS', 'AUTHORIZED REPRESENTATIVES',
    'authorized signers', 'authorized representatives',
}

# ---------------------------------------------------------------------------
# Extraction
# ---------------------------------------------------------------------------

def esc(s):
    return s.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')


def is_noise(text):
    """Check if text is repeated noise (headers/footers on every slide)."""
    stripped = text.strip()
    return stripped in NOISE_TEXTS


def extract_slide_structured(slide):
    """Extract slide content with layout awareness.
    Returns dict with: title, subtitle, main_texts, sidebar_texts,
    tables, images, notes, source_texts
    """
    result = {
        'title': '',
        'subtitle': '',
        'main_texts': [],
        'sidebar_texts': [],
        'tables': [],
        'images': [],  # list of (base64_data, content_type, width_in, height_in, position)
        'notes': '',
        'source_texts': [],
    }

    for shape in slide.shapes:
        left = shape.left or 0
        top = shape.top or 0
        is_sidebar = left > SIDEBAR_THRESHOLD
        is_title_area = top < TITLE_THRESHOLD
        is_source_area = top > SOURCE_THRESHOLD

        # Images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                w_in = shape.width / EMU_PER_INCH
                h_in = shape.height / EMU_PER_INCH
                # Skip tiny icons (< 0.5 inch)
                if w_in < 0.5 and h_in < 0.5:
                    continue
                b64 = base64.b64encode(img.blob).decode('ascii')
                ct = img.content_type
                pos = 'sidebar' if is_sidebar else 'main'
                result['images'].append((b64, ct, w_in, h_in, pos))
            except Exception:
                pass
            continue

        # Tables
        if shape.has_table:
            t = shape.table
            rows = []
            for row in t.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            result['tables'].append(rows)
            continue

        # Text shapes
        if not hasattr(shape, 'text') or not shape.text.strip():
            continue

        text = shape.text.strip()

        # Filter noise
        if is_noise(text):
            continue

        # Classify by position
        if is_source_area and (text.lower().startswith('source') or len(text) < 200):
            # Could be source attribution at bottom
            if text.lower().startswith('source'):
                result['source_texts'].append(text)
                continue

        if is_title_area and not result['title'] and len(text) < 120:
            result['title'] = text
        elif is_title_area and result['title'] and not result['subtitle'] and len(text) < 200:
            result['subtitle'] = text
        elif is_sidebar:
            result['sidebar_texts'].append(text)
        else:
            result['main_texts'].append(text)

    # Notes
    if slide.has_notes_slide:
        result['notes'] = slide.notes_slide.notes_text_frame.text.strip()

    return result


# ---------------------------------------------------------------------------
# HTML formatting
# ---------------------------------------------------------------------------

def format_text(text):
    """Smart text formatting with bullet detection, key:value, etc."""
    lines = text.split('\n')
    result = []
    in_list = False
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        # Bullet/dash items
        if re.match(r'^[•·–\-]\s', stripped) or re.match(r'^[•·]\s*', stripped) and len(stripped) > 2:
            if not in_list:
                result.append('<ul class="styled-list">')
                in_list = True
            item = re.sub(r'^[•·–\-]\s*', '', stripped)
            result.append(f'<li>{esc(item)}</li>')
        elif re.match(r'^\d+[\.\)]\s', stripped):
            if not in_list:
                result.append('<ul class="styled-list">')
                in_list = True
            item = re.sub(r'^\d+[\.\)]\s*', '', stripped)
            result.append(f'<li>{esc(item)}</li>')
        else:
            if in_list:
                result.append('</ul>')
                in_list = False
            # ALL CAPS heading
            if stripped.isupper() and 4 < len(stripped) < 80:
                result.append(f'<h4 class="content-heading">{esc(stripped.title())}</h4>')
            # Key: Value pattern
            elif re.match(r'^[A-Z][^:]{1,35}:\s+\S', stripped):
                parts = stripped.split(':', 1)
                result.append(f'<p class="kv-line"><strong>{esc(parts[0])}:</strong> {esc(parts[1].strip())}</p>')
            else:
                result.append(f'<p class="body-text">{esc(stripped)}</p>')

    if in_list:
        result.append('</ul>')
    return '\n'.join(result)


def table_to_html(rows):
    if not rows:
        return ''
    html = ['<div class="table-wrap"><table>']
    html.append('<thead><tr>')
    for cell in rows[0]:
        html.append(f'<th>{esc(cell)}</th>')
    html.append('</tr></thead><tbody>')
    for row in rows[1:]:
        html.append('<tr>')
        for cell in row:
            html.append(f'<td>{esc(cell)}</td>')
        html.append('</tr>')
    html.append('</tbody></table></div>')
    return '\n'.join(html)


def images_to_html(images, position_filter=None):
    """Render images as embedded HTML."""
    html = []
    for b64, ct, w_in, h_in, pos in images:
        if position_filter and pos != position_filter:
            continue
        # Scale images to fit: max width ~100%, maintain aspect
        max_w = 100 if pos == 'main' else 100
        html.append(f'<div class="slide-image"><img src="data:{ct};base64,{b64}" style="max-width:{max_w}%; height:auto; border-radius:6px; border:1px solid var(--border);" alt="Slide visual"></div>')
    return '\n'.join(html)


def render_slide(slide_data, slide_num):
    """Render a single slide's content as HTML."""
    d = slide_data
    parts = []

    has_sidebar = bool(d['sidebar_texts'])
    has_main_images = any(1 for _,_,_,_,p in d['images'] if p == 'main')
    has_sidebar_images = any(1 for _,_,_,_,p in d['images'] if p == 'sidebar')

    # Main content
    main_parts = []
    for txt in d['main_texts']:
        main_parts.append(format_text(txt))
    for tbl in d['tables']:
        main_parts.append(table_to_html(tbl))
    main_img_html = images_to_html(d['images'], 'main')
    if main_img_html:
        main_parts.append(main_img_html)

    # Sidebar content
    sidebar_parts = []
    for txt in d['sidebar_texts']:
        sidebar_parts.append(format_text(txt))
    sidebar_img_html = images_to_html(d['images'], 'sidebar')
    if sidebar_img_html:
        sidebar_parts.append(sidebar_img_html)

    main_html = '\n'.join(main_parts) if main_parts else ''
    sidebar_html = '\n'.join(sidebar_parts) if sidebar_parts else ''

    # Notes
    notes_html = ''
    if d['notes']:
        notes_html = f'''<div class="notes-box">
    <h4>Notes</h4>
    <div>{format_text(d["notes"])}</div>
</div>'''

    # Source
    source_html = ''
    if d['source_texts']:
        sources = ' · '.join(d['source_texts'])
        source_html = f'<div class="source-tag">{esc(sources)}</div>'

    # Build the card
    title_text = d['title'] if d['title'] else f'Slide {slide_num}'
    subtitle_html = f'<div class="slide-subtitle">{esc(d["subtitle"])}</div>' if d['subtitle'] else ''

    if not main_html and not sidebar_html and not d['images']:
        body_inner = '<p class="empty-note">Visual content — refer to original presentation</p>'
    elif has_sidebar and (sidebar_html):
        body_inner = f'''{subtitle_html}
<div class="two-col">
    <div>{main_html}</div>
    <div class="sidebar-panel">{sidebar_html}</div>
</div>'''
    else:
        body_inner = f'''{subtitle_html}
{main_html}'''

    return f'''<div class="card">
    <div class="card-header">
        <span>{esc(title_text)}</span>
        <span class="slide-num">Slide {slide_num}</span>
    </div>
    <div class="card-body">
        {body_inner}
        {notes_html}
        {source_html}
    </div>
</div>'''


# ---------------------------------------------------------------------------
# Section builder
# ---------------------------------------------------------------------------

def build_section(slides, slide_nums, section_title, kpi_bar_html=''):
    """Build a complete section from slide numbers."""
    parts = []
    if kpi_bar_html:
        parts.append(kpi_bar_html)

    for sn in slide_nums:
        if sn > len(slides):
            continue
        data = extract_slide_structured(slides[sn - 1])
        # Skip completely empty slides
        if not data['title'] and not data['main_texts'] and not data['sidebar_texts'] \
           and not data['tables'] and not data['images']:
            continue
        parts.append(render_slide(data, sn))

    return '\n'.join(parts) if parts else '<p class="empty-note">No extractable content for this section.</p>'


def exec_stat(label, value, sub='', color='var(--primary)'):
    sub_html = f'<div class="sub">{esc(sub)}</div>' if sub else ''
    return f'''<div class="exec-stat">
    <div class="label">{esc(label)}</div>
    <div class="value" style="color:{color};">{esc(str(value))}</div>
    {sub_html}
</div>'''


# ---------------------------------------------------------------------------
# CSS
# ---------------------------------------------------------------------------

CSS = '''
:root {
    --primary: #1a365d;
    --primary-light: #2c5282;
    --bg: #f7fafc;
    --card-bg: #ffffff;
    --text: #2d3748;
    --text-light: #718096;
    --border: #e2e8f0;
    --green: #276749;
    --green-bg: #f0fff4;
    --green-border: #c6f6d5;
    --red: #c53030;
    --red-bg: #fff5f5;
    --red-border: #fed7d7;
    --amber: #b7791f;
    --amber-bg: #fffff0;
    --amber-border: #fefcbf;
    --blue-bg: #ebf8ff;
    --blue-border: #bee3f8;
}
* { margin:0; padding:0; box-sizing:border-box; }
body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; background:var(--bg); color:var(--text); font-size:14px; line-height:1.6; }

/* Header */
.header { background:linear-gradient(135deg, var(--primary), var(--primary-light)); color:#fff; padding:28px 32px; }
.header h1 { font-size:22px; font-weight:700; }
.header .sub { font-size:13px; opacity:0.85; margin-top:4px; }

.container { max-width:1440px; margin:0 auto; padding:20px; }

/* Executive Stats */
.exec-bar { display:grid; grid-template-columns:repeat(4, 1fr); gap:16px; margin-bottom:24px; }
.exec-stat { background:var(--card-bg); border-radius:8px; padding:18px; box-shadow:0 1px 3px rgba(0,0,0,0.08); text-align:center; }
.exec-stat .label { font-size:11px; text-transform:uppercase; letter-spacing:0.5px; color:var(--text-light); margin-bottom:4px; }
.exec-stat .value { font-size:28px; font-weight:700; }
.exec-stat .sub { font-size:11px; color:var(--text-light); margin-top:2px; }

/* Tabs */
.tabs { display:flex; gap:2px; background:var(--border); border-radius:8px 8px 0 0; padding:4px 4px 0; overflow-x:auto; flex-wrap:nowrap; }
.tab { padding:10px 14px; cursor:pointer; background:#e2e8f0; border-radius:6px 6px 0 0; font-size:12px; font-weight:500; white-space:nowrap; border:none; color:var(--text); }
.tab:hover { background:#cbd5e0; }
.tab.active { background:var(--card-bg); color:var(--primary); font-weight:600; }
.tab-content { display:none; padding:20px 0; }
.tab-content.active { display:block; }

/* Cards */
.card { background:var(--card-bg); border-radius:8px; box-shadow:0 1px 3px rgba(0,0,0,0.08); margin-bottom:20px; overflow:hidden; }
.card-header { background:var(--bg); padding:12px 16px; font-weight:600; font-size:14px; border-bottom:1px solid var(--border); display:flex; justify-content:space-between; align-items:center; }
.card-body { padding:16px; }
.slide-num { font-size:11px; color:var(--text-light); font-weight:400; background:var(--border); padding:2px 8px; border-radius:10px; }
.slide-subtitle { font-size:13px; color:var(--text-light); margin-bottom:12px; font-style:italic; }

/* Two column layout */
.two-col { display:grid; grid-template-columns:1fr 1fr; gap:20px; }
.sidebar-panel { background:var(--bg); border-radius:6px; padding:14px; border:1px solid var(--border); }

/* Typography */
.content-heading { font-size:14px; font-weight:600; color:var(--primary); margin:14px 0 8px; padding-bottom:4px; border-bottom:2px solid var(--primary); }
.body-text { font-size:13px; margin-bottom:8px; line-height:1.7; }
.kv-line { font-size:13px; margin-bottom:6px; }
.empty-note { color:var(--text-light); font-size:12px; font-style:italic; }

/* Lists */
.styled-list { list-style:none; padding:0; margin:8px 0; }
.styled-list li { padding:6px 12px; border-bottom:1px solid var(--border); font-size:13px; line-height:1.6; position:relative; padding-left:24px; }
.styled-list li:last-child { border-bottom:none; }
.styled-list li::before { content:"•"; color:var(--primary); font-weight:700; position:absolute; left:8px; }

/* Images */
.slide-image { margin:12px 0; text-align:center; }
.slide-image img { display:inline-block; }

/* Notes */
.notes-box { background:var(--amber-bg); border:1px solid var(--amber-border); border-radius:6px; padding:12px 16px; margin-top:14px; }
.notes-box h4 { color:var(--amber); font-size:11px; text-transform:uppercase; letter-spacing:0.5px; margin-bottom:6px; }
.notes-box p, .notes-box li { font-size:12px; color:#744210; }

/* Source */
.source-tag { font-size:10px; color:var(--text-light); font-style:italic; margin-top:10px; padding-top:8px; border-top:1px solid var(--border); }

/* Tables */
.table-wrap { overflow-x:auto; margin:12px 0; }
table { width:100%; border-collapse:collapse; font-size:13px; }
th { background:var(--primary); color:#fff; padding:8px 12px; text-align:left; font-size:12px; }
td { padding:8px 12px; border-bottom:1px solid var(--border); }
tr:hover { background:#edf2f7; }

/* Info & Alert Boxes */
.info-box { background:var(--blue-bg); border:1px solid var(--blue-border); border-radius:6px; padding:14px 18px; margin-bottom:14px; }
.info-box h4 { color:#2b6cb0; font-size:13px; margin-bottom:6px; }

/* Responsive */
@media (max-width:900px) {
    .two-col { grid-template-columns:1fr; }
    .exec-bar { grid-template-columns:1fr 1fr; }
    .tabs { flex-wrap:wrap; }
}
'''


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    base = os.path.dirname(os.path.abspath(__file__))
    pptx_path = os.path.join(base, 'GMA Current State.pptx')

    print(f"Reading {pptx_path}...")
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    print(f"Found {len(slides)} slides")

    # Section definitions
    section_defs = [
        ('regulatory', 'Regulatory Framework', [2,3,4,5,6,7,8],
         f'''<div class="exec-bar" style="grid-template-columns:repeat(3,1fr);">
{exec_stat('Region', 'WAEMU', '8-member monetary union', 'var(--primary)')}
{exec_stat('Central Bank', 'BCEAO', "Banque Centrale des États de l'Afrique de l'Ouest")}
{exec_stat('Currency', 'XOF', 'Pegged to EUR (655.957)')}
</div>'''),
        ('payments', 'Payment Systems', [9,10,11,12,13,14,15],
         f'''<div class="exec-bar">
{exec_stat('RTGS', 'STAR-UEMOA', 'High-value settlement', 'var(--green)')}
{exec_stat('ACH', 'SICA-UEMOA', 'Bulk clearing')}
{exec_stat('Interbank', 'GIM-UEMOA', 'Card & mobile', 'var(--amber)')}
{exec_stat('New', 'PI-SPI', 'Payment initiation', 'var(--primary-light)')}
</div>'''),
        ('structure', 'Business & Legal', [16,17,18,19,20],
         f'''<div class="exec-bar">
{exec_stat('Entity', 'GMA', "Grands Moulins d'Abidjan", 'var(--primary)')}
{exec_stat('Country', "Côte d'Ivoire", 'WAEMU member')}
{exec_stat('Industry', 'Flour Milling', 'Wheat processing & distribution')}
{exec_stat('Est.', '1963', 'Post-independence')}
</div>'''),
        ('tax', 'Tax Overview', [21], ''),
        ('repatriation', 'Funding & Repatriation', [22,23,24,25], ''),
        ('technology', 'Technology', [26], ''),
        ('accounts', 'Bank Accounts', [27,28,29,30,31,32,33,34,35,36],
         f'''<div class="exec-bar" style="grid-template-columns:repeat(3,1fr);">
{exec_stat('Accounts', '10', 'Across 8 banks', 'var(--primary)')}
{exec_stat('Currency', 'XOF', 'West African CFA Franc')}
{exec_stat('Coverage', '2024–2025', 'Transaction activity data', 'var(--green)')}
</div>'''),
        ('operations', 'Operations', [37,38,39,40,41], ''),
        ('signers', 'Users & Signers', [42], ''),
        ('fx', 'FX & Trade', [43,44,45,46], ''),
        ('facilities', 'Facilities', [47,48], ''),
        ('pativoire', 'Pativoire', [49,50,51],
         f'''<div class="exec-bar" style="grid-template-columns:repeat(3,1fr);">
{exec_stat('Entity', 'Pativoire', 'Subsidiary / affiliate')}
{exec_stat('Country', "Côte d'Ivoire", 'Same jurisdiction as GMA')}
{exec_stat('Status', 'Active', 'Opening Citi account', 'var(--amber)')}
</div>'''),
        ('painpoints', 'Pain Points', [52], ''),
        ('fundsavail', 'Funds Availability', [53,54,55,56], ''),
        ('questions', 'Questions', [57,58,59,60], ''),
    ]

    sections = {}
    total_images = 0
    for key, title, slide_nums, kpi_html in section_defs:
        print(f"  Building {key} (slides {slide_nums[0]}-{slide_nums[-1]})...")
        content = build_section(slides, slide_nums, title, kpi_html)
        # Count embedded images
        for sn in slide_nums:
            if sn <= len(slides):
                for shape in slides[sn-1].shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        w = shape.width / EMU_PER_INCH
                        h = shape.height / EMU_PER_INCH
                        if w >= 0.5 or h >= 0.5:
                            total_images += 1
        sections[key] = {'title': title, 'content': content}

    # Build tab UI
    tab_buttons = []
    tab_contents = []
    first = True
    for key, title, _, _ in section_defs:
        active = ' active' if first else ''
        tab_buttons.append(f'<button class="tab{active}" onclick="showTab_cs(\'{key}\')">{sections[key]["title"]}</button>')
        tab_contents.append(f'<div id="{key}" class="tab-content{active}">\n{sections[key]["content"]}\n</div>')
        first = False

    html = f'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>GMA Current State — From PowerPoint ({len(slides)} Slides)</title>
<style>
{CSS}
</style>
</head>
<body>
<div class="header">
    <h1>GMA — Les Grands Moulins d'Abidjan</h1>
    <div class="sub">Current State | Côte d'Ivoire | From PowerPoint ({len(slides)} Slides, {total_images} visuals embedded) | March 2026</div>
</div>
<div class="container">

<div class="tabs">
    {"".join(tab_buttons)}
</div>

{"".join(tab_contents)}

</div>
<script>
function showTab_cs(id) {{
    document.querySelectorAll('.tab-content').forEach(t => t.classList.remove('active'));
    document.querySelectorAll('.tab').forEach(t => t.classList.remove('active'));
    document.getElementById(id).classList.add('active');
    event.target.closest('.tab').classList.add('active');
}}
</script>
</body>
</html>'''

    out_path = os.path.join(base, 'gma-current-state-ppt.html')
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(html)

    size_kb = os.path.getsize(out_path) / 1024
    print(f"\nDone! Output: {out_path} ({size_kb:.0f} KB)")
    print(f"Sections: {len(sections)}, Embedded images: {total_images}")
    for key, title, nums, _ in section_defs:
        print(f"  {key}: {title} (slides {nums[0]}-{nums[-1]})")


if __name__ == '__main__':
    main()
